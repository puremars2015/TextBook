from pptx import Presentation

import os
import shutil

# 打開PPT檔案
prs = Presentation('C:\\Users\\sean.ma\\Downloads\\週報20260109馬恩奇.pptx')

# 刪除第7、8、9張投影片（注意：使用反向索引刪除以避免索引變化）
# 要刪除的投影片索引：6, 7, 8 (因為索引從0開始)
for idx in [8, 7, 6]:  # 反向刪除
    rId = prs.slides._sldIdLst[idx].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[idx]

print(f"刪除後投影片數: {len(prs.slides)}")

# 先保存到臨時檔案
temp_file = 'C:\\Users\\sean.ma\\Downloads\\週報20260109馬恩奇_temp.pptx'
prs.save(temp_file)
print("臨時檔案已保存")

# 用臨時檔案替換原始檔案
import time
time.sleep(2)  # 等待檔案釋放
# 先嘗試刪除
try:
    os.remove('C:\\Users\\sean.ma\\Downloads\\週報20260109馬恩奇.pptx')
    shutil.move(temp_file, 'C:\\Users\\sean.ma\\Downloads\\週報20260109馬恩奇.pptx')
    print("PPT檔案已修復")
except PermissionError:
    # 如果無法刪除，嘗試覆蓋
    shutil.move(temp_file, 'C:\\Users\\sean.ma\\Downloads\\週報20260109馬恩奇.pptx', copy_function=shutil.copy2)
    print("PPT檔案已用覆蓋方式修復")
