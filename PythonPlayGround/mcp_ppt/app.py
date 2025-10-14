#!/usr/bin/env python3
"""
MCP Server for PowerPoint Generation
依照 README.md 規格實作，使用 FastMCP + python-pptx
"""

import os
import re
import logging
from datetime import datetime
from pathlib import Path
from typing import Optional, Dict, Any, List
from uuid import uuid4

from pptx import Presentation
from pptx.util import Inches, Pt
from fastmcp import FastMCP

# ==================== 設定常數 ====================
OUTPUT_DIR = "./output"
MAX_SESSIONS = 8
DEFAULT_FILENAME_FMT = "ppt_{timestamp}.pptx"

# ==================== 初始化 ====================
# 建立輸出資料夾
os.makedirs(OUTPUT_DIR, exist_ok=True)

# 設定日誌
log_file = os.path.join(OUTPUT_DIR, "mcp_ppt.log")
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler(log_file),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

# 建立 FastMCP 實例
mcp = FastMCP("PPT Server")

# ==================== 全域狀態 ====================
# 儲存樣板與工作中的簡報
template_prs: Optional[Presentation] = None
template_path: Optional[str] = None
sessions: Dict[str, Presentation] = {}

# ==================== 工具函數 ====================

def get_absolute_path(path: str) -> str:
    """取得標準化的絕對路徑"""
    return str(Path(path).resolve())


def get_layout_by_name_or_index(prs: Presentation, layout_identifier: str | int) -> tuple:
    """
    根據名稱或索引取得版面配置
    回傳 (layout, warning)
    """
    layouts = prs.slide_layouts
    
    # 嘗試作為索引
    if isinstance(layout_identifier, int):
        try:
            return layouts[layout_identifier], None
        except IndexError:
            logger.warning(f"Layout index {layout_identifier} out of range, using layout 0")
            return layouts[0], f"Layout index {layout_identifier} not found, used default"
    
    # 嘗試作為名稱
    if isinstance(layout_identifier, str):
        # 先嘗試精確匹配
        for i, layout in enumerate(layouts):
            if layout.name == layout_identifier:
                return layout, None
        
        # 嘗試作為數字字串
        try:
            idx = int(layout_identifier)
            if 0 <= idx < len(layouts):
                return layouts[idx], None
        except ValueError:
            pass
        
        # 失敗，使用預設
        logger.warning(f"Layout '{layout_identifier}' not found, using layout 0")
        return layouts[0], f"Layout '{layout_identifier}' not found, used default"
    
    return layouts[0], None


def parse_description_dsl(text: str, assets_base_dir: Optional[str] = None) -> tuple:
    """
    解析文字敘述 DSL
    回傳 (metadata, slides, warnings)
    """
    warnings = []
    metadata = {}
    slides = []
    
    # 解析開頭的 metadata
    lines = text.strip().split('\n')
    i = 0
    
    # 解析 metadata (# 開頭)
    while i < len(lines):
        line = lines[i].strip()
        if line.startswith('#'):
            if ':' in line:
                key, value = line[1:].split(':', 1)
                key = key.strip().lower()
                value = value.strip()
                if key in ['簡報標題', 'title']:
                    metadata['title'] = value
                elif key in ['作者', 'author']:
                    metadata['author'] = value
                elif key in ['日期', 'date']:
                    metadata['date'] = value
            i += 1
        elif not line:
            i += 1
        else:
            break
    
    # 解析投影片
    current_slide = None
    current_bullets = []
    in_bullets = False
    
    while i < len(lines):
        line = lines[i].strip()
        
        if line == '[Slide]':
            # 儲存前一張投影片
            if current_slide is not None:
                if current_bullets:
                    current_slide['bullets'] = current_bullets
                slides.append(current_slide)
            
            # 開始新投影片
            current_slide = {
                'layout': None,
                'title': None,
                'subtitle': None,
                'bullets': None,
                'image': None,
                'notes': None
            }
            current_bullets = []
            in_bullets = False
            
        elif current_slide is not None and ':' in line:
            key, value = line.split(':', 1)
            key = key.strip().lower()
            value = value.strip()
            
            if key == 'layout':
                current_slide['layout'] = value
            elif key == 'title':
                current_slide['title'] = value
            elif key == 'subtitle':
                current_slide['subtitle'] = value
            elif key == 'bullets':
                in_bullets = True
            elif key == 'image':
                # 處理圖片路徑
                img_path = value
                if assets_base_dir and not os.path.isabs(img_path):
                    img_path = os.path.join(assets_base_dir, img_path)
                current_slide['image'] = img_path
            elif key == 'notes':
                current_slide['notes'] = value
                in_bullets = False
        elif line.startswith('-') and current_slide is not None:
            # 項目符號
            bullet_text = line[1:].strip()
            current_bullets.append(bullet_text)
        elif not line:
            in_bullets = False
        
        i += 1
    
    # 儲存最後一張投影片
    if current_slide is not None:
        if current_bullets:
            current_slide['bullets'] = current_bullets
        slides.append(current_slide)
    
    return metadata, slides, warnings


def add_slide_content(slide, layout, title: Optional[str], subtitle: Optional[str],
                      bullets: Optional[List[str]], notes: Optional[str]) -> List[str]:
    """
    填充投影片內容
    回傳 warnings
    """
    warnings = []
    
    try:
        # 填充標題
        if title and hasattr(slide.shapes, 'title') and slide.shapes.title:
            slide.shapes.title.text = title
        
        # 填充副標題
        if subtitle:
            # 嘗試找到副標題佔位符
            for shape in slide.shapes:
                if hasattr(shape, 'placeholder_format'):
                    ph_type = shape.placeholder_format.type
                    if ph_type == 2:  # PP_PLACEHOLDER.SUBTITLE
                        shape.text = subtitle
                        break
        
        # 填充項目符號
        if bullets:
            # 尋找內容佔位符
            for shape in slide.shapes:
                if hasattr(shape, 'placeholder_format'):
                    ph_type = shape.placeholder_format.type
                    if ph_type == 2 or ph_type == 7:  # BODY or OBJECT
                        if shape.has_text_frame:
                            text_frame = shape.text_frame
                            text_frame.clear()
                            for i, bullet in enumerate(bullets):
                                if i == 0:
                                    p = text_frame.paragraphs[0]
                                else:
                                    p = text_frame.add_paragraph()
                                p.text = bullet
                                p.level = 0
                            break
        
        # 填充備註
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes
            
    except Exception as e:
        warning = f"Error adding content: {str(e)}"
        warnings.append(warning)
        logger.warning(warning)
    
    return warnings


# ==================== MCP 工具 ====================

@mcp.tool()
def set_template(template_path_input: str) -> Dict[str, Any]:
    """
    設定樣板簡報，並回傳可用版面配置
    
    Args:
        template_path_input: 現有 .pptx 路徑
    """
    global template_prs, template_path
    
    try:
        abs_path = get_absolute_path(template_path_input)
        
        if not os.path.exists(abs_path):
            return {
                "ok": False,
                "error_code": "TEMPLATE_NOT_FOUND",
                "message": f"Template file not found: {abs_path}"
            }
        
        # 載入樣板
        template_prs = Presentation(abs_path)
        template_path = abs_path
        
        # 取得版面配置資訊
        layouts = []
        for i, layout in enumerate(template_prs.slide_layouts):
            layouts.append({
                "index": i,
                "name": layout.name
            })
        
        # 嘗試取得主題資訊（有限支援）
        theme = {
            "color_scheme_names": [],
            "font_families": []
        }
        
        logger.info(f"Template set: {abs_path}")
        
        return {
            "ok": True,
            "template_abs_path": abs_path,
            "layouts": layouts,
            "theme": theme
        }
        
    except Exception as e:
        logger.error(f"Error setting template: {str(e)}")
        return {
            "ok": False,
            "error_code": "FILE_IO_ERROR",
            "message": str(e)
        }


@mcp.tool()
def create_presentation(use_template: bool = True, 
                       metadata: Optional[Dict[str, str]] = None) -> Dict[str, Any]:
    """
    建立工作中的簡報物件
    
    Args:
        use_template: 是否使用樣板
        metadata: 包含 title, author, date 的字典
    """
    global sessions, template_prs
    
    try:
        if len(sessions) >= MAX_SESSIONS:
            return {
                "ok": False,
                "error_code": "UNSUPPORTED",
                "message": f"Maximum sessions ({MAX_SESSIONS}) reached"
            }
        
        # 建立新簡報
        if use_template and template_prs:
            prs = Presentation(template_path)
            # 清空所有投影片
            while len(prs.slides) > 0:
                rId = prs.slides._sldIdLst[0].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[0]
        else:
            prs = Presentation()
            if use_template and not template_prs:
                logger.warning("Template requested but not set, using blank presentation")
        
        # 設定 metadata
        if metadata:
            core_props = prs.core_properties
            if 'title' in metadata and metadata['title']:
                core_props.title = metadata['title']
            if 'author' in metadata and metadata['author']:
                core_props.author = metadata['author']
            if 'date' in metadata and metadata['date']:
                core_props.comments = f"Date: {metadata['date']}"
        
        # 產生 session_id
        session_id = str(uuid4())
        sessions[session_id] = prs
        
        logger.info(f"Presentation created: {session_id}")
        
        return {
            "ok": True,
            "session_id": session_id
        }
        
    except Exception as e:
        logger.error(f"Error creating presentation: {str(e)}")
        return {
            "ok": False,
            "error_code": "PPTX_WRITE_ERROR",
            "message": str(e)
        }


@mcp.tool()
def add_slide(session_id: str, 
             layout: str | int,
             title: Optional[str] = None,
             subtitle: Optional[str] = None,
             bullets: Optional[List[str]] = None,
             image: Optional[str] = None,
             notes: Optional[str] = None) -> Dict[str, Any]:
    """
    以結構化資料新增單張投影片
    
    Args:
        session_id: 工作階段 ID
        layout: 版面配置名稱或索引
        title: 標題
        subtitle: 副標題
        bullets: 項目符號列表
        image: 圖片路徑
        notes: 講者備註
    """
    try:
        if session_id not in sessions:
            return {
                "ok": False,
                "error_code": "SESSION_NOT_FOUND",
                "message": f"Session not found: {session_id}"
            }
        
        prs = sessions[session_id]
        warnings = []
        
        # 取得版面配置
        slide_layout, warning = get_layout_by_name_or_index(prs, layout)
        if warning:
            warnings.append(warning)
        
        # 新增投影片
        slide = prs.slides.add_slide(slide_layout)
        slide_index = len(prs.slides) - 1
        
        # 填充內容
        content_warnings = add_slide_content(slide, slide_layout, title, subtitle, bullets, notes)
        warnings.extend(content_warnings)
        
        # 插入圖片
        if image:
            try:
                if os.path.exists(image):
                    # 自動置中並縮放
                    left = Inches(1)
                    top = Inches(2)
                    slide.shapes.add_picture(image, left, top, width=Inches(8), height=Inches(5))
                else:
                    warning = f"Image not found: {image}"
                    warnings.append(warning)
                    logger.warning(warning)
            except Exception as e:
                warning = f"Error inserting image: {str(e)}"
                warnings.append(warning)
                logger.warning(warning)
        
        logger.info(f"Slide added to session {session_id}, index: {slide_index}")
        
        result = {
            "ok": True,
            "index": slide_index
        }
        if warnings:
            result["warnings"] = warnings
        
        return result
        
    except Exception as e:
        logger.error(f"Error adding slide: {str(e)}")
        return {
            "ok": False,
            "error_code": "PPTX_WRITE_ERROR",
            "message": str(e)
        }


@mcp.tool()
def compile_from_description(description_text: str,
                            template_path_input: Optional[str] = None,
                            assets_base_dir: Optional[str] = None) -> Dict[str, Any]:
    """
    解析「文字敘述 DSL」，批次建立簡報
    
    Args:
        description_text: DSL 文字
        template_path_input: 樣板路徑（可選）
        assets_base_dir: 素材資料夾路徑
    """
    global template_prs, template_path
    
    try:
        warnings = []
        
        # 設定樣板
        if template_path_input:
            result = set_template(template_path_input)
            if not result.get("ok"):
                return result
        
        # 解析 DSL
        metadata, slides_data, parse_warnings = parse_description_dsl(description_text, assets_base_dir)
        warnings.extend(parse_warnings)
        
        # 建立簡報
        create_result = create_presentation(
            use_template=bool(template_prs),
            metadata=metadata if metadata else None
        )
        
        if not create_result.get("ok"):
            return create_result
        
        session_id = create_result["session_id"]
        
        # 新增投影片
        for slide_data in slides_data:
            if not slide_data.get('layout'):
                slide_data['layout'] = 0
            
            add_result = add_slide(
                session_id=session_id,
                layout=slide_data['layout'],
                title=slide_data.get('title'),
                subtitle=slide_data.get('subtitle'),
                bullets=slide_data.get('bullets'),
                image=slide_data.get('image'),
                notes=slide_data.get('notes')
            )
            
            if add_result.get("ok") and add_result.get("warnings"):
                warnings.extend(add_result["warnings"])
        
        logger.info(f"Compiled presentation from description, session: {session_id}")
        
        return {
            "ok": True,
            "session_id": session_id,
            "slide_count": len(slides_data),
            "warnings": warnings
        }
        
    except Exception as e:
        logger.error(f"Error compiling from description: {str(e)}")
        return {
            "ok": False,
            "error_code": "PARSE_ERROR",
            "message": str(e)
        }


@mcp.tool()
def list_layouts() -> Dict[str, Any]:
    """
    列出目前樣板或工作中的簡報可用版面
    """
    try:
        prs = template_prs if template_prs else Presentation()
        
        layouts = []
        for i, layout in enumerate(prs.slide_layouts):
            layouts.append({
                "index": i,
                "name": layout.name
            })
        
        return {
            "ok": True,
            "layouts": layouts
        }
        
    except Exception as e:
        logger.error(f"Error listing layouts: {str(e)}")
        return {
            "ok": False,
            "error_code": "PPTX_WRITE_ERROR",
            "message": str(e)
        }


@mcp.tool()
def insert_image(session_id: str,
                slide_index: int,
                image_path: str,
                left: Optional[float] = None,
                top: Optional[float] = None,
                width: Optional[float] = None,
                height: Optional[float] = None) -> Dict[str, Any]:
    """
    在指定投影片插圖
    
    Args:
        session_id: 工作階段 ID
        slide_index: 投影片索引
        image_path: 圖片路徑
        left: 左邊距（英吋）
        top: 上邊距（英吋）
        width: 寬度（英吋）
        height: 高度（英吋）
    """
    try:
        if session_id not in sessions:
            return {
                "ok": False,
                "error_code": "SESSION_NOT_FOUND",
                "message": f"Session not found: {session_id}"
            }
        
        prs = sessions[session_id]
        
        if slide_index < 0 or slide_index >= len(prs.slides):
            return {
                "ok": False,
                "error_code": "INVALID_LAYOUT",
                "message": f"Slide index {slide_index} out of range"
            }
        
        if not os.path.exists(image_path):
            return {
                "ok": False,
                "error_code": "IMAGE_NOT_FOUND",
                "message": f"Image not found: {image_path}"
            }
        
        slide = prs.slides[slide_index]
        
        # 處理尺寸參數
        kwargs = {}
        if left is not None:
            kwargs['left'] = Inches(left)
        else:
            kwargs['left'] = Inches(1)
            
        if top is not None:
            kwargs['top'] = Inches(top)
        else:
            kwargs['top'] = Inches(2)
            
        if width is not None:
            kwargs['width'] = Inches(width)
        if height is not None:
            kwargs['height'] = Inches(height)
        
        # 如果沒有指定寬高，使用預設
        if 'width' not in kwargs and 'height' not in kwargs:
            kwargs['width'] = Inches(8)
        
        slide.shapes.add_picture(image_path, **kwargs)
        
        logger.info(f"Image inserted in session {session_id}, slide {slide_index}")
        
        return {
            "ok": True
        }
        
    except Exception as e:
        logger.error(f"Error inserting image: {str(e)}")
        return {
            "ok": False,
            "error_code": "IMAGE_NOT_FOUND",
            "message": str(e)
        }


@mcp.tool()
def finalize_and_save(session_id: str,
                     output_dir: Optional[str] = None,
                     filename: Optional[str] = None) -> Dict[str, Any]:
    """
    儲存並結束工作階段
    
    Args:
        session_id: 工作階段 ID
        output_dir: 輸出資料夾
        filename: 檔案名稱
    """
    try:
        if session_id not in sessions:
            return {
                "ok": False,
                "error_code": "SESSION_NOT_FOUND",
                "message": f"Session not found: {session_id}"
            }
        
        prs = sessions[session_id]
        
        # 決定輸出路徑
        out_dir = output_dir if output_dir else OUTPUT_DIR
        os.makedirs(out_dir, exist_ok=True)
        
        if not filename:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = DEFAULT_FILENAME_FMT.replace("{timestamp}", timestamp)
        
        if not filename.endswith('.pptx'):
            filename += '.pptx'
        
        output_path = os.path.join(out_dir, filename)
        abs_path = get_absolute_path(output_path)
        
        # 儲存
        prs.save(abs_path)
        
        # 取得檔案大小
        size_bytes = os.path.getsize(abs_path)
        
        # 清理 session
        del sessions[session_id]
        
        logger.info(f"Presentation saved: {abs_path}")
        
        return {
            "ok": True,
            "pptx_abs_path": abs_path,
            "size_bytes": size_bytes
        }
        
    except Exception as e:
        logger.error(f"Error saving presentation: {str(e)}")
        return {
            "ok": False,
            "error_code": "PPTX_WRITE_ERROR",
            "message": str(e)
        }


@mcp.tool()
def reset() -> Dict[str, Any]:
    """
    釋放所有工作階段與樣板快取
    """
    global sessions, template_prs, template_path
    
    try:
        sessions.clear()
        template_prs = None
        template_path = None
        
        logger.info("All sessions and template cache cleared")
        
        return {
            "ok": True
        }
        
    except Exception as e:
        logger.error(f"Error resetting: {str(e)}")
        return {
            "ok": False,
            "error_code": "UNSUPPORTED",
            "message": str(e)
        }


# ==================== 主程式 ====================

if __name__ == "__main__":
    import sys
    
    # 預設使用 streamable HTTP 模式
    mode = "http"
    port = 8000
    host = "127.0.0.1"
    
    # 解析命令列參數
    if len(sys.argv) > 1:
        if sys.argv[1] == "stdio":
            mode = "stdio"
        elif sys.argv[1] == "sse":
            mode = "sse"
        elif sys.argv[1] == "http":
            mode = "http"
    
    # 檢查是否有指定 port
    if len(sys.argv) > 2:
        try:
            port = int(sys.argv[2])
        except ValueError:
            logger.warning(f"Invalid port number: {sys.argv[2]}, using default {port}")
    
    # 啟動 server
    if mode == "stdio":
        logger.info("Starting MCP PPT Server in stdio mode...")
        mcp.run()
    elif mode == "sse":
        logger.info(f"Starting MCP PPT Server in SSE mode on {host}:{port}...")
        import uvicorn
        uvicorn.run(mcp.sse_app, host=host, port=port)
    else:  # http (streamable)
        logger.info(f"Starting MCP PPT Server in streamable HTTP mode on {host}:{port}...")
        import uvicorn
        app = mcp.streamable_http_app()
        uvicorn.run(app, host=host, port=port)
