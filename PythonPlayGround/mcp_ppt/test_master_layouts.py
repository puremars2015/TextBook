#!/usr/bin/env python3
"""
深入測試投影片母片和版面配置
"""

import os
from pptx import Presentation
from pptx.oxml import parse_xml

def test_slide_master_info():
    """檢查投影片母片的詳細資訊"""
    template_path = "/Users/maenqi/TextBook/PythonPlayGround/mcp_ppt/template.pptx"
    
    print("=" * 80)
    print("深入分析投影片母片 (Slide Master)")
    print("=" * 80)
    
    if not os.path.exists(template_path):
        print(f"❌ 找不到檔案: {template_path}")
        return
    
    prs = Presentation(template_path)
    
    print(f"📊 簡報資訊:")
    print(f"  - 投影片數量: {len(prs.slides)}")
    print(f"  - 投影片母片數量: {len(prs.slide_masters)}")
    print()
    
    # 檢查每個投影片母片
    for master_idx, slide_master in enumerate(prs.slide_masters):
        print(f"投影片母片 #{master_idx + 1}")
        print("-" * 80)
        print(f"  名稱: {slide_master.name if hasattr(slide_master, 'name') else 'N/A'}")
        print(f"  版面配置數量: {len(slide_master.slide_layouts)}")
        print()
        
        print("  版面配置清單:")
        for layout_idx, layout in enumerate(slide_master.slide_layouts):
            print(f"    [{layout_idx}] {layout.name}")
            
            # 檢查版面配置的佔位符
            placeholders = []
            for shape in layout.shapes:
                if hasattr(shape, 'placeholder_format'):
                    ph_type = shape.placeholder_format.type
                    ph_idx = shape.placeholder_format.idx
                    placeholders.append(f"Type={ph_type}, Idx={ph_idx}")
            
            if placeholders:
                print(f"         佔位符: {', '.join(placeholders)}")
        print()
    
    print()
    print("=" * 80)
    print("透過 prs.slide_layouts 直接存取")
    print("=" * 80)
    print(f"總共可用的版面配置數量: {len(prs.slide_layouts)}")
    print()
    for i, layout in enumerate(prs.slide_layouts):
        print(f"  [{i}] {layout.name}")
    print()
    
    # 檢查實際投影片使用的版面
    print()
    print("=" * 80)
    print("現有投影片使用的版面配置")
    print("=" * 80)
    for slide_idx, slide in enumerate(prs.slides):
        layout = slide.slide_layout
        print(f"投影片 #{slide_idx + 1}: 使用版面 '{layout.name}'")
    print()


def test_compare_with_blank():
    """與空白簡報比較"""
    print()
    print("=" * 80)
    print("對照：空白簡報的母片結構")
    print("=" * 80)
    
    blank_prs = Presentation()
    
    print(f"空白簡報的投影片母片數量: {len(blank_prs.slide_masters)}")
    
    for master_idx, slide_master in enumerate(blank_prs.slide_masters):
        print(f"\n母片 #{master_idx + 1}:")
        print(f"  版面配置數量: {len(slide_master.slide_layouts)}")
        for layout_idx, layout in enumerate(slide_master.slide_layouts):
            print(f"    [{layout_idx}] {layout.name}")


def test_xml_structure():
    """檢查 XML 結構（進階）"""
    template_path = "/Users/maenqi/TextBook/PythonPlayGround/mcp_ppt/template.pptx"
    
    print()
    print("=" * 80)
    print("檢查內部 XML 結構")
    print("=" * 80)
    
    if not os.path.exists(template_path):
        print(f"❌ 找不到檔案: {template_path}")
        return
    
    try:
        prs = Presentation(template_path)
        
        # 檢查簡報的 part 關係
        print("簡報包含的 parts:")
        for rel in prs.part.rels.values():
            print(f"  - {rel.reltype}: {rel.target_part}")
        
        print()
        
        # 檢查母片的關係
        if len(prs.slide_masters) > 0:
            master = prs.slide_masters[0]
            print(f"\n第一個母片的 parts:")
            for rel in master.part.rels.values():
                print(f"  - {rel.reltype}")
                if 'slideLayout' in rel.reltype:
                    print(f"    -> Layout part: {rel.target_part}")
        
    except Exception as e:
        print(f"⚠️  無法完整解析 XML: {str(e)}")


if __name__ == "__main__":
    print("\n" + "🔍 深入分析投影片母片與版面配置 🔍".center(80))
    print("\n")
    
    test_slide_master_info()
    test_compare_with_blank()
    test_xml_structure()
    
    print("\n" + "=" * 80)
    print("✅ 分析完成")
    print("=" * 80)
    
    print("\n💡 說明:")
    print("   如果您的 template.pptx 只顯示 1 種版面配置，")
    print("   可能是因為該範本檔案在建立時，只保留了一個投影片母片版面。")
    print("   ")
    print("   建議解決方案:")
    print("   1. 在 PowerPoint 中開啟 template.pptx")
    print("   2. 進入「檢視」→「投影片母片」")
    print("   3. 確認左側有多個版面配置（應該要有 6 個）")
    print("   4. 如果只有一個，請從其他範本複製版面配置過來")
    print("   5. 儲存後再測試")
