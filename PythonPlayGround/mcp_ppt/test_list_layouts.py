#!/usr/bin/env python3
"""
測試 list_layouts 功能
檢查範本檔案中有多少種版型可用
"""

import os
from pptx import Presentation

def test_list_layouts_directly():
    """直接使用 python-pptx 測試範本的版型數量"""
    template_path = "/Users/maenqi/TextBook/PythonPlayGround/mcp_ppt/template.pptx"
    
    print("=" * 60)
    print("測試 1: 直接載入範本檔案")
    print("=" * 60)
    
    if not os.path.exists(template_path):
        print(f"❌ 錯誤: 找不到範本檔案: {template_path}")
        return
    
    print(f"✅ 範本檔案存在: {template_path}")
    print()
    
    # 載入範本
    prs = Presentation(template_path)
    
    print(f"📊 範本中的投影片數量: {len(prs.slides)}")
    print(f"📐 範本中的版面配置數量: {len(prs.slide_layouts)}")
    print()
    
    print("版面配置清單:")
    print("-" * 60)
    for i, layout in enumerate(prs.slide_layouts):
        print(f"  索引 {i}: {layout.name}")
    print()


def test_list_layouts_after_clear():
    """測試清空投影片後，版型是否還在"""
    template_path = "/Users/maenqi/TextBook/PythonPlayGround/mcp_ppt/template.pptx"
    
    print("=" * 60)
    print("測試 2: 清空投影片後的版型數量")
    print("=" * 60)
    
    if not os.path.exists(template_path):
        print(f"❌ 錯誤: 找不到範本檔案: {template_path}")
        return
    
    # 載入範本
    prs = Presentation(template_path)
    print(f"清空前 - 投影片數量: {len(prs.slides)}")
    print(f"清空前 - 版面配置數量: {len(prs.slide_layouts)}")
    print()
    
    # 使用原本的方法清空投影片 (while 版本)
    print("使用 while 迴圈清空投影片...")
    prs1 = Presentation(template_path)
    while len(prs1.slides) > 0:
        rId = prs1.slides._sldIdLst[0].rId
        prs1.part.drop_rel(rId)
        del prs1.slides._sldIdLst[0]
    
    print(f"清空後 (while) - 投影片數量: {len(prs1.slides)}")
    print(f"清空後 (while) - 版面配置數量: {len(prs1.slide_layouts)}")
    print("版面配置清單:")
    for i, layout in enumerate(prs1.slide_layouts):
        print(f"  索引 {i}: {layout.name}")
    print()
    
    # 使用修改後的方法清空投影片 (for 版本)
    print("使用 for 迴圈清空投影片...")
    prs2 = Presentation(template_path)
    for i in range(len(prs2.slides) - 1, -1, -1):
        rId = prs2.slides._sldIdLst[i].rId
        prs2.part.drop_rel(rId)
        del prs2.slides._sldIdLst[i]
    
    print(f"清空後 (for) - 投影片數量: {len(prs2.slides)}")
    print(f"清空後 (for) - 版面配置數量: {len(prs2.slide_layouts)}")
    print("版面配置清單:")
    for i, layout in enumerate(prs2.slide_layouts):
        print(f"  索引 {i}: {layout.name}")
    print()


def test_list_layouts_with_app_logic():
    """測試使用 app.py 中的邏輯"""
    template_path = "/Users/maenqi/TextBook/PythonPlayGround/mcp_ppt/template.pptx"
    
    print("=" * 60)
    print("測試 3: 模擬 app.py 的 list_layouts 邏輯")
    print("=" * 60)
    
    if not os.path.exists(template_path):
        print(f"❌ 錯誤: 找不到範本檔案: {template_path}")
        return
    
    # 模擬 set_template
    print("步驟 1: 設定樣板 (set_template)")
    template_prs = Presentation(template_path)
    print(f"✅ 樣板已載入")
    print(f"📐 版面配置數量: {len(template_prs.slide_layouts)}")
    print()
    
    # 模擬 create_presentation
    print("步驟 2: 建立簡報 (create_presentation)")
    prs = Presentation(template_path)
    print(f"清空前 - 投影片數量: {len(prs.slides)}")
    
    # 清空投影片
    for i in range(len(prs.slides) - 1, -1, -1):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]
    
    print(f"清空後 - 投影片數量: {len(prs.slides)}")
    print(f"清空後 - 版面配置數量: {len(prs.slide_layouts)}")
    print()
    
    # 模擬 list_layouts (使用 session)
    print("步驟 3: 列出版面配置 (list_layouts with session)")
    layouts = []
    for i, layout in enumerate(prs.slide_layouts):
        layouts.append({
            "index": i,
            "name": layout.name
        })
    
    print(f"📊 找到 {len(layouts)} 種版面配置:")
    for layout in layouts:
        print(f"  索引 {layout['index']}: {layout['name']}")
    print()


def test_blank_presentation():
    """測試空白簡報的版型數量（作為對照）"""
    print("=" * 60)
    print("測試 4: 空白簡報的版型數量（對照組）")
    print("=" * 60)
    
    prs = Presentation()
    print(f"📐 空白簡報的版面配置數量: {len(prs.slide_layouts)}")
    print()
    print("版面配置清單:")
    print("-" * 60)
    for i, layout in enumerate(prs.slide_layouts):
        print(f"  索引 {i}: {layout.name}")
    print()


if __name__ == "__main__":
    print("\n" + "🔍 開始測試 list_layouts 功能 🔍".center(60))
    print("\n")
    
    # 執行所有測試
    test_list_layouts_directly()
    test_list_layouts_after_clear()
    test_list_layouts_with_app_logic()
    test_blank_presentation()
    
    print("=" * 60)
    print("✅ 測試完成")
    print("=" * 60)
