from fastmcp import FastMCP
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
from typing import Optional

# 建立 FastMCP 伺服器實例
mcp = FastMCP("PowerPoint MCP Server")

# 全域變數存儲當前打開的簡報
current_presentation: Optional[Presentation] = None
current_file_path: Optional[str] = None

# 定義主題相關的顏色方案
THEME_COLORS = {
    "商業": {
        "primary": RGBColor(0, 51, 102),      # 深藍色
        "secondary": RGBColor(51, 102, 153),  # 藍色
        "accent": RGBColor(255, 153, 0),      # 橙色
        "background": RGBColor(255, 255, 255) # 白色
    },
    "科技": {
        "primary": RGBColor(0, 150, 136),     # 青綠色
        "secondary": RGBColor(0, 188, 212),   # 天藍色
        "accent": RGBColor(255, 193, 7),      # 金黃色
        "background": RGBColor(250, 250, 250) # 淺灰白
    },
    "教育": {
        "primary": RGBColor(63, 81, 181),     # 靛藍色
        "secondary": RGBColor(103, 58, 183),  # 深紫色
        "accent": RGBColor(233, 30, 99),      # 粉紅色
        "background": RGBColor(255, 255, 255) # 白色
    },
    "醫療": {
        "primary": RGBColor(0, 121, 107),     # 深青色
        "secondary": RGBColor(38, 166, 154),  # 青色
        "accent": RGBColor(255, 87, 34),      # 橘紅色
        "background": RGBColor(255, 255, 255) # 白色
    },
    "創意": {
        "primary": RGBColor(156, 39, 176),    # 紫色
        "secondary": RGBColor(233, 30, 99),   # 粉紅色
        "accent": RGBColor(255, 193, 7),      # 黃色
        "background": RGBColor(255, 255, 255) # 白色
    },
    "預設": {
        "primary": RGBColor(68, 68, 68),      # 深灰色
        "secondary": RGBColor(128, 128, 128), # 灰色
        "accent": RGBColor(0, 112, 192),      # 藍色
        "background": RGBColor(255, 255, 255) # 白色
    }
}


@mcp.tool()
def open_ppt(file_path: str) -> str:
    """
    開啟現有的 PowerPoint 檔案
    
    Args:
        file_path: PPT 檔案的路徑
    
    Returns:
        操作結果訊息
    """
    global current_presentation, current_file_path
    
    try:
        if not os.path.exists(file_path):
            return f"錯誤: 檔案 '{file_path}' 不存在"
        
        current_presentation = Presentation(file_path)
        current_file_path = file_path
        
        slide_count = len(current_presentation.slides)
        return f"成功開啟檔案: {file_path}\n總共有 {slide_count} 張投影片"
    
    except Exception as e:
        return f"開啟檔案時發生錯誤: {str(e)}"


@mcp.tool()
def create_ppt(file_path: str, theme: Optional[str] = None, use_template: bool = False) -> str:
    """
    創建新的 PowerPoint 檔案並套用佈景主題
    
    Args:
        file_path: 新 PPT 檔案的保存路徑
        theme: 主題名稱（可選：商業、科技、教育、醫療、創意，若不提供則自動根據檔名選擇）
        use_template: 是否使用 template_v2.pptx 作為模板（預設為 False）
    
    Returns:
        操作結果訊息
    """
    global current_presentation, current_file_path
    
    try:
        if use_template:
            # 使用模板創建
            template_path = os.path.join(os.path.dirname(__file__), "template_v2.pptx")
            if not os.path.exists(template_path):
                return f"錯誤: 模板檔案 'template_v2.pptx' 不存在"
            
            current_presentation = Presentation(template_path)
            current_file_path = file_path
            
            # 保存為新檔案
            current_presentation.save(file_path)
            
            slide_count = len(current_presentation.slides)
            layout_count = len(current_presentation.slide_layouts)
            return f"成功使用模板創建 PowerPoint 檔案: {file_path}\n模板已包含 {slide_count} 張投影片，{layout_count} 種版型"
        else:
            # 創建空白簡報
            current_presentation = Presentation()
            current_file_path = file_path
            
            # 自動選擇主題
            if theme is None:
                theme = _auto_select_theme(file_path)
            
            # 套用主題顏色
            applied_theme = _apply_theme(theme)
            
            # 保存新建的簡報
            current_presentation.save(file_path)
            
            return f"成功創建新的 PowerPoint 檔案: {file_path}\n套用的佈景主題: {applied_theme}"
    
    except Exception as e:
        return f"創建檔案時發生錯誤: {str(e)}"


@mcp.tool()
def create_from_template(file_path: str, template_path: str = "template_v2.pptx") -> str:
    """
    使用指定的模板檔案創建新的 PowerPoint 檔案
    
    Args:
        file_path: 新 PPT 檔案的保存路徑
        template_path: 模板檔案路徑（預設為 template_v2.pptx）
    
    Returns:
        操作結果訊息
    """
    global current_presentation, current_file_path
    
    try:
        # 如果是相對路徑，在當前目錄尋找
        if not os.path.isabs(template_path):
            template_path = os.path.join(os.path.dirname(__file__), template_path)
        
        if not os.path.exists(template_path):
            return f"錯誤: 模板檔案 '{template_path}' 不存在"
        
        # 使用模板創建新簡報
        current_presentation = Presentation(template_path)
        current_file_path = file_path
        
        # 保存為新檔案
        current_presentation.save(file_path)
        
        slide_count = len(current_presentation.slides)
        layout_count = len(current_presentation.slide_layouts)
        
        return f"成功使用模板創建 PowerPoint 檔案: {file_path}\n模板: {os.path.basename(template_path)}\n包含 {slide_count} 張投影片，{layout_count} 種版型"
    
    except Exception as e:
        return f"使用模板創建檔案時發生錯誤: {str(e)}"


def _auto_select_theme(file_path: str) -> str:
    """
    根據檔案名稱自動選擇合適的主題
    
    Args:
        file_path: 檔案路徑
    
    Returns:
        主題名稱
    """
    file_name = os.path.basename(file_path).lower()
    
    # 根據關鍵字匹配主題
    theme_keywords = {
        "商業": ["商業", "business", "企業", "公司", "finance", "財務", "報告", "report"],
        "科技": ["科技", "tech", "technology", "ai", "軟體", "software", "程式", "code", "開發", "dev"],
        "教育": ["教育", "education", "學校", "school", "教學", "teaching", "課程", "course", "學習", "learn"],
        "醫療": ["醫療", "medical", "health", "健康", "醫院", "hospital", "診所", "clinic"],
        "創意": ["創意", "creative", "設計", "design", "藝術", "art", "品牌", "brand"]
    }
    
    for theme_name, keywords in theme_keywords.items():
        for keyword in keywords:
            if keyword in file_name:
                return theme_name
    
    return "預設"


def _apply_theme(theme_name: str) -> str:
    """
    套用主題顏色配置到簡報
    
    Args:
        theme_name: 主題名稱
    
    Returns:
        實際套用的主題名稱
    """
    global current_presentation
    
    # 如果主題名稱不存在，使用預設主題
    if theme_name not in THEME_COLORS:
        theme_name = "預設"
    
    colors = THEME_COLORS[theme_name]
    
    # 套用主題到簡報的母片
    try:
        # 獲取簡報的母片
        for slide_master in current_presentation.slide_master:
            # 設定背景顏色
            background = slide_master.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = colors["background"]
            
            # 設定標題顏色
            for shape in slide_master.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = colors["primary"]
    except Exception as e:
        # 如果無法套用到母片，至少記錄主題選擇
        pass
    
    return theme_name


@mcp.tool()
def write_ppt(
    title: Optional[str] = None,
    content: Optional[str] = None,
    layout_index: Optional[int] = None
) -> str:
    """
    在 PowerPoint 檔案中寫入新的投影片
    
    Args:
        title: 投影片標題
        content: 投影片內容
        layout_index: 使用的版型索引（若不提供則自動選擇）
    
    Returns:
        操作結果訊息
    """
    global current_presentation, current_file_path
    
    if current_presentation is None:
        return "錯誤: 尚未開啟或創建 PowerPoint 檔案"
    
    try:
        # 自動選擇版型
        if layout_index is None:
            layout_index = _auto_select_layout(title, content)
        
        # 選擇版型
        slide_layout = current_presentation.slide_layouts[layout_index]
        
        # 新增投影片
        slide = current_presentation.slides.add_slide(slide_layout)
        
        # 設定標題
        if title and slide.shapes.title:
            slide.shapes.title.text = title
        
        # 設定內容
        if content:
            # 尋找內容文字框
            content_set = False
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 2:  # 內容占位符
                    shape.text = content
                    content_set = True
                    break
            
            # 如果沒有找到內容占位符，嘗試使用文字框
            if not content_set:
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and shape != slide.shapes.title:
                        shape.text_frame.text = content
                        content_set = True
                        break
        
        # 保存檔案
        if current_file_path:
            current_presentation.save(current_file_path)
        
        layout_name = slide_layout.name
        slide_num = len(current_presentation.slides)
        return f"成功新增投影片（第 {slide_num} 張），使用版型: {layout_name} (索引 {layout_index})"
    
    except Exception as e:
        return f"寫入投影片時發生錯誤: {str(e)}"


def _auto_select_layout(title: Optional[str], content: Optional[str]) -> int:
    """
    根據內容自動選擇合適的版型
    
    Args:
        title: 標題
        content: 內容
    
    Returns:
        版型索引
    """
    global current_presentation
    
    # 取得所有版型
    layouts = current_presentation.slide_layouts
    
    # 如果只有標題沒有內容
    if title and not content:
        # 尋找標題投影片或區段標題
        for idx, layout in enumerate(layouts):
            layout_name = layout.name.lower()
            if 'title' in layout_name and 'content' not in layout_name:
                return idx
        return 0  # 預設使用第一個版型
    
    # 如果有標題和內容
    if title and content:
        # 尋找標題+內容版型
        for idx, layout in enumerate(layouts):
            layout_name = layout.name.lower()
            if 'title' in layout_name and 'content' in layout_name:
                return idx
        # 尋找包含 "Title and Content" 或類似的版型
        for idx, layout in enumerate(layouts):
            if idx == 1:  # 通常索引1是標題+內容版型
                return idx
        return 1 if len(layouts) > 1 else 0
    
    # 如果只有內容沒有標題
    if content and not title:
        # 尋找空白或純內容版型
        for idx, layout in enumerate(layouts):
            layout_name = layout.name.lower()
            if 'blank' in layout_name or 'content' in layout_name:
                return idx
        return 6 if len(layouts) > 6 else 0
    
    # 都沒有，返回空白版型
    for idx, layout in enumerate(layouts):
        if 'blank' in layout.name.lower():
            return idx
    
    return 0  # 預設第一個


@mcp.tool()
def edit_ppt(
    slide_index: int,
    title: Optional[str] = None,
    content: Optional[str] = None
) -> str:
    """
    編輯現有的投影片
    
    Args:
        slide_index: 要編輯的投影片索引（從0開始）
        title: 新的標題（若不提供則不修改）
        content: 新的內容（若不提供則不修改）
    
    Returns:
        操作結果訊息
    """
    global current_presentation, current_file_path
    
    if current_presentation is None:
        return "錯誤: 尚未開啟或創建 PowerPoint 檔案"
    
    try:
        # 檢查投影片索引是否有效
        if slide_index >= len(current_presentation.slides):
            return f"錯誤: 投影片索引 {slide_index} 超出範圍（總共 {len(current_presentation.slides)} 張）"
        
        slide = current_presentation.slides[slide_index]
        
        # 更新標題
        if title and slide.shapes.title:
            slide.shapes.title.text = title
        
        # 更新內容
        if content:
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 2:  # 內容占位符
                    shape.text = content
                    break
        
        # 保存檔案
        if current_file_path:
            current_presentation.save(current_file_path)
        
        return f"成功編輯投影片 {slide_index}"
    
    except Exception as e:
        return f"編輯投影片時發生錯誤: {str(e)}"


@mcp.tool()
def list_layouts() -> str:
    """
    列出當前 PowerPoint 檔案中可用的所有版型
    
    Returns:
        版型列表及其索引
    """
    global current_presentation
    
    if current_presentation is None:
        return "錯誤: 尚未開啟或創建 PowerPoint 檔案"
    
    try:
        layouts_info = []
        for idx, layout in enumerate(current_presentation.slide_layouts):
            layout_name = layout.name
            layouts_info.append(f"索引 {idx}: {layout_name}")
        
        result = "可用的投影片版型:\n" + "\n".join(layouts_info)
        return result
    
    except Exception as e:
        return f"列出版型時發生錯誤: {str(e)}"


@mcp.tool()
def get_ppt_info() -> str:
    """
    取得當前 PowerPoint 檔案的資訊
    
    Returns:
        檔案資訊
    """
    global current_presentation, current_file_path
    
    if current_presentation is None:
        return "目前沒有開啟的 PowerPoint 檔案"
    
    try:
        slide_count = len(current_presentation.slides)
        layout_count = len(current_presentation.slide_layouts)
        
        info = f"""
當前檔案資訊:
- 檔案路徑: {current_file_path}
- 投影片數量: {slide_count}
- 可用版型數量: {layout_count}
        """
        return info.strip()
    
    except Exception as e:
        return f"取得資訊時發生錯誤: {str(e)}"


# 使用 HTTP streamable 協定啟動伺服器
if __name__ == "__main__":
    # 使用 HTTP 協定
    mcp.run(transport="http", host="127.0.0.1", port=8700)