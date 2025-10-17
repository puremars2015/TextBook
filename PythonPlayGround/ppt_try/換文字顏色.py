from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

# 設定你要統一的顏色
TARGET_COLOR = RGBColor(47, 85, 151)  # 深藍

def recolor_all_shapes(src, dst):
    prs = Presentation(src)
    count = 0

    for slide in prs.slides:
        for shape in slide.shapes:
            # 跳過群組形狀
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                continue

            # 嘗試填滿顏色（僅限自訂形狀）
            if shape.shape_type in (
                MSO_SHAPE_TYPE.AUTO_SHAPE,
                MSO_SHAPE_TYPE.FREEFORM,
                MSO_SHAPE_TYPE.PICTURE,
                MSO_SHAPE_TYPE.TEXT_BOX,
            ):
                try:
                    fill = shape.fill
                    fill.solid()
                    fill.fore_color.rgb = TARGET_COLOR
                    count += 1
                except Exception:
                    pass

            # 修改文字顏色（若有文字框）
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        font = run.font
                        font.color.rgb = TARGET_COLOR

    prs.save(dst)
    print(f"已輸出 {dst}，共修改 {count} 個圖形顏色。")

# 執行
if __name__ == "__main__":
    recolor_all_shapes("input.pptx", "output_recolor.pptx")
