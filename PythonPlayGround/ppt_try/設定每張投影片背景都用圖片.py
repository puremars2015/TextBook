from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

# === 顏色設定 ===
TARGET_COLOR = RGBColor(47, 85, 151)  # 圖形與文字顏色（深藍）
BACKGROUND_IMAGE = "background.png"   # 你的背景圖檔路徑

def recolor_all_shapes_with_bgimage(src, dst):
    prs = Presentation(src)
    shape_count = 0
    slide_count = 0

    for slide in prs.slides:
        # --- 設定底圖 ---
        slide.shapes.add_picture(
            BACKGROUND_IMAGE,
            Inches(0), Inches(0),
            prs.slide_width, prs.slide_height
        )
        slide_count += 1

        # --- 修改每個 shape ---
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                continue

            # 填滿顏色（可視物件）
            if shape.shape_type in (
                MSO_SHAPE_TYPE.AUTO_SHAPE,
                MSO_SHAPE_TYPE.FREEFORM,
                MSO_SHAPE_TYPE.PICTURE,
                MSO_SHAPE_TYPE.TEXT_BOX,
            ):
                try:
                    fill = shape.fill
                    fill.solid()
                    fill.fore_color.rgb = TARGET_COLOR
                    shape_count += 1
                except Exception:
                    pass

            # 文字顏色
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        font = run.font
                        font.color.rgb = TARGET_COLOR

    prs.save(dst)
    print(f"已輸出 {dst}，共 {slide_count} 張投影片、修改 {shape_count} 個圖形。")

# === 執行 ===
if __name__ == "__main__":
    recolor_all_shapes_with_bgimage("input.pptx", "output_bgimg.pptx")
