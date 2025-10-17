from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

# 統一顏色設定
TARGET_COLOR = RGBColor(47, 85, 151)   # shape與文字顏色（深藍）
BACKGROUND_COLOR = RGBColor(230, 240, 255)  # 投影片底色（淡藍）

def recolor_all_shapes_with_bg(src, dst):
    prs = Presentation(src)
    shape_count = 0
    slide_count = 0

    for slide in prs.slides:
        # 設定底色
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = BACKGROUND_COLOR
        slide_count += 1

        for shape in slide.shapes:
            # 跳過群組
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                continue

            # 改填滿顏色
            if shape.shape_type in (
                MSO_SHAPE_TYPE.AUTO_SHAPE,
                MSO_SHAPE_TYPE.FREEFORM,
                MSO_SHAPE_TYPE.PICTURE,
                MSO_SHAPE_TYPE.TEXT_BOX,
            ):
                try:
                    fill = shape.fill
                    fill.solid()
                    fill.fore_color.rgb = TARGET_COLOR
                    shape_count += 1
                except Exception:
                    pass

            # 改文字顏色
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        font = run.font
                        font.color.rgb = TARGET_COLOR

    prs.save(dst)
    print(f"已輸出 {dst}，修改 {slide_count} 張投影片背景，{shape_count} 個圖形。")

# 執行
if __name__ == "__main__":
    recolor_all_shapes_with_bg("input.pptx", "output_recolor_bg.pptx")
