from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_pad_course_presentation():
    """
    創建Power Automate Desktop課程的PowerPoint簡報
    """
    # 創建新的簡報
    prs = Presentation()
    
    # 設定簡報尺寸 (16:9)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # 定義顏色主題
    primary_color = RGBColor(0, 120, 212)  # 微軟藍
    secondary_color = RGBColor(40, 40, 40)  # 深灰
    accent_color = RGBColor(255, 185, 0)   # 橘黃
    
    # 1. 標題頁
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版面
    
    # 添加標題
    title_shape = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(1.5))
    title_frame = title_shape.text_frame
    title_frame.text = "Power Automate Desktop\n實務入門與進階"
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(48)
    title_p.font.bold = True
    title_p.font.color.rgb = primary_color
    title_p.alignment = PP_ALIGN.CENTER
    
    # 添加副標題
    subtitle_shape = slide1.shapes.add_textbox(Inches(1), Inches(4), Inches(11.33), Inches(1))
    subtitle_frame = subtitle_shape.text_frame
    subtitle_frame.text = "RPA自動化解決方案完整教學"
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.font.size = Pt(24)
    subtitle_p.font.color.rgb = secondary_color
    subtitle_p.alignment = PP_ALIGN.CENTER
    
    # 2. 課程目標
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 標題
    add_slide_title(slide2, "課程目標", primary_color)
    
    # 內容
    objectives = [
        "讓學員能夠理解 RPA（機器人流程自動化）基本概念與 Power Automate Desktop (PAD) 的角色",
        "掌握 PAD 的介面與常用動作，能夠設計、建置、偵錯並部署桌面自動化流程",
        "透過實作練習，把 PAD 與 Excel、Outlook、Web 瀏覽器等常見應用整合",
        "提供進階主題（API、UiPath/Power Platform整合、效能與例外處理）供實務應用"
    ]
    
    add_bullet_points(slide2, objectives, Inches(1), Inches(2.5))
    
    # 3. 目標受眾
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_slide_title(slide3, "目標受眾", primary_color)
    
    # 受眾類型
    audience_title_shape = slide3.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(0.5))
    audience_title_frame = audience_title_shape.text_frame
    audience_title_frame.text = "適合對象："
    audience_title_p = audience_title_frame.paragraphs[0]
    audience_title_p.font.size = Pt(20)
    audience_title_p.font.bold = True
    audience_title_p.font.color.rgb = accent_color
    
    audience_list = [
        "業務流程優化專員",
        "IT 支援人員", 
        "自動化開發初學者",
        "產品經理",
        "任何希望提高工作效率的上班族"
    ]
    
    add_bullet_points(slide3, audience_list, Inches(1.5), Inches(3.2))
    
    # 先備知識
    prereq_title_shape = slide3.shapes.add_textbox(Inches(1), Inches(5), Inches(11.33), Inches(0.5))
    prereq_title_frame = prereq_title_shape.text_frame
    prereq_title_frame.text = "先備知識："
    prereq_title_p = prereq_title_frame.paragraphs[0]
    prereq_title_p.font.size = Pt(20)
    prereq_title_p.font.bold = True
    prereq_title_p.font.color.rgb = accent_color
    
    prereq_text_shape = slide3.shapes.add_textbox(Inches(1.5), Inches(5.7), Inches(10), Inches(1))
    prereq_text_frame = prereq_text_shape.text_frame
    prereq_text_frame.text = "具備基本電腦操作能力與 Office 軟體使用經驗；若有程式基礎（如 VBA、Python）將更容易學習進階內容"
    prereq_text_p = prereq_text_frame.paragraphs[0]
    prereq_text_p.font.size = Pt(16)
    prereq_text_p.font.color.rgb = secondary_color
    
    # 4. 教學方式與時數建議
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_slide_title(slide4, "教學方式與時數建議", primary_color)
    
    teaching_info = [
        "授課形式：講解 + 示範 + 分組實作 + 作業回饋",
        "建議時數：24 小時（3 日密集班）或 8 週每週 3 小時（夜間班）",
        "每章節搭配練習與小專案",
        "理論/實作比例：約 30% 理論 / 70% 實作"
    ]
    
    add_bullet_points(slide4, teaching_info, Inches(1), Inches(2.5))
    
    # 5. 課程大綱 - 第一章
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_slide_title(slide5, "課程大綱 - 第一章", primary_color)
    
    # 章節標題
    chapter1_title_shape = slide5.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.33), Inches(0.5))
    chapter1_title_frame = chapter1_title_shape.text_frame
    chapter1_title_frame.text = "1. 認識 RPA 與 Power Automate Desktop（90 分鐘）"
    chapter1_title_p = chapter1_title_frame.paragraphs[0]
    chapter1_title_p.font.size = Pt(22)
    chapter1_title_p.font.bold = True
    chapter1_title_p.font.color.rgb = accent_color
    
    # 目的
    add_section_content(slide5, "目的：", "建立自動化思維，了解 PAD 在微軟生態系的定位", Inches(1), Inches(2.9))
    
    # 要點
    chapter1_points = [
        "RPA 概念介紹",
        "PAD 與 Power Automate（雲端）差異",
        "授權與安全性考量",
        "應用場景示例"
    ]
    add_section_with_bullets(slide5, "要點：", chapter1_points, Inches(1), Inches(3.6))
    
    # 實作
    add_section_content(slide5, "實作：", "安裝 PAD 與基礎設定、建立第一個簡單流程（Hello world：開啟應用程式 + 訊息視窗）", Inches(1), Inches(5.5))
    
    # 6. 課程大綱 - 第二章
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_slide_title(slide6, "課程大綱 - 第二章", primary_color)
    
    # 章節標題
    chapter2_title_shape = slide6.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.33), Inches(0.5))
    chapter2_title_frame = chapter2_title_shape.text_frame
    chapter2_title_frame.text = "2. PAD 介面導覽與流程基礎（120 分鐘）"
    chapter2_title_p = chapter2_title_frame.paragraphs[0]
    chapter2_title_p.font.size = Pt(22)
    chapter2_title_p.font.bold = True
    chapter2_title_p.font.color.rgb = accent_color
    
    # 目的
    add_section_content(slide6, "目的：", "熟悉 PAD Studio 的工作區、動作庫、變數與流程結構", Inches(1), Inches(2.9))
    
    # 要點
    chapter2_points = [
        "動作分類（UI、自動化、資料、流程控制）",
        "錄製（Recorder）功能",
        "變數管理",
        "流程分段與子流程"
    ]
    add_section_with_bullets(slide6, "要點：", chapter2_points, Inches(1), Inches(3.6))
    
    # 實作
    add_section_content(slide6, "實作：", "錄製使用者操作、編輯動作、觀察變數、建立子流程", Inches(1), Inches(5.5))
    
    # 7. 課程大綱 - 第三章
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_slide_title(slide7, "課程大綱 - 第三章", primary_color)
    
    # 章節標題
    chapter3_title_shape = slide7.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.33), Inches(0.5))
    chapter3_title_frame = chapter3_title_shape.text_frame
    chapter3_title_frame.text = "3. 控制流程與錯誤處理（120 分鐘）"
    chapter3_title_p = chapter3_title_frame.paragraphs[0]
    chapter3_title_p.font.size = Pt(22)
    chapter3_title_p.font.bold = True
    chapter3_title_p.font.color.rgb = accent_color
    
    # 目的
    add_section_content(slide7, "目的：", "掌握條件判斷、迴圈、例外處理與日誌", Inches(1), Inches(2.9))
    
    # 要點
    chapter3_points = [
        "If/Else 條件判斷",
        "For Each、While 迴圈",
        "Try/Catch 例外處理",
        "Log 記錄",
        "調試技巧（斷點、逐步執行）"
    ]
    add_section_with_bullets(slide7, "要點：", chapter3_points, Inches(1), Inches(3.6))
    
    # 實作
    add_section_content(slide7, "實作：", "建立含條件判斷與迴圈的流程，模擬並處理錯誤情境", Inches(1), Inches(5.8))
    
    # 儲存簡報
    prs.save('Power_Automate_Desktop_Course.pptx')
    print("PowerPoint 簡報已成功創建：Power_Automate_Desktop_Course.pptx")

def add_slide_title(slide, title_text, color):
    """添加投影片標題"""
    title_shape = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.33), Inches(1))
    title_frame = title_shape.text_frame
    title_frame.text = title_text
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(36)
    title_p.font.bold = True
    title_p.font.color.rgb = color
    title_p.alignment = PP_ALIGN.LEFT

def add_bullet_points(slide, points, left, top):
    """添加項目符號列表"""
    content_shape = slide.shapes.add_textbox(left, top, Inches(10.5), Inches(4))
    content_frame = content_shape.text_frame
    
    for i, point in enumerate(points):
        if i == 0:
            p = content_frame.paragraphs[0]
        else:
            p = content_frame.add_paragraph()
        
        p.text = point
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(40, 40, 40)
        p.level = 0

def add_section_content(slide, label, content, left, top):
    """添加區段內容（標籤 + 內容）"""
    section_shape = slide.shapes.add_textbox(left, top, Inches(11), Inches(0.6))
    section_frame = section_shape.text_frame
    section_frame.text = label + content
    section_p = section_frame.paragraphs[0]
    section_p.font.size = Pt(14)
    
    # 設定標籤樣式
    run = section_p.runs[0]
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 185, 0)  # 橘黃色

def add_section_with_bullets(slide, label, points, left, top):
    """添加區段標籤和項目符號列表"""
    # 標籤
    label_shape = slide.shapes.add_textbox(left, top, Inches(11), Inches(0.4))
    label_frame = label_shape.text_frame
    label_frame.text = label
    label_p = label_frame.paragraphs[0]
    label_p.font.size = Pt(14)
    label_p.font.bold = True
    label_p.font.color.rgb = RGBColor(255, 185, 0)
    
    # 項目符號
    bullets_shape = slide.shapes.add_textbox(Inches(left.inches + 0.5), Inches(top.inches + 0.4), Inches(10), Inches(1.5))
    bullets_frame = bullets_shape.text_frame
    
    for i, point in enumerate(points):
        if i == 0:
            p = bullets_frame.paragraphs[0]
        else:
            p = bullets_frame.add_paragraph()
        
        p.text = point
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(40, 40, 40)
        p.level = 0

if __name__ == "__main__":
    try:
        create_pad_course_presentation()
    except ImportError:
        print("請先安裝 python-pptx 套件：")
        print("pip install python-pptx")
    except Exception as e:
        print(f"發生錯誤：{e}")