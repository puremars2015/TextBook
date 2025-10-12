"""
AI-Powered PPT Generator
類似 Gamma 的自動簡報生成服務
使用 Claude API 自動生成專業簡報
"""

import os
import json
from typing import Any, Dict, List, Optional
from dataclasses import dataclass
from datetime import datetime
import anthropic
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER

@dataclass
class SlideContent:
    """簡報頁面內容結構"""
    title: str
    content: List[str]
    layout: str = "content"  # title, content, two_column, image_text
    notes: str = ""

class AIpptGenerator:
    """AI 驅動的 PPT 生成器"""
    
    def __init__(self, api_key: str):
        """
        初始化生成器
        
        Args:
            api_key: Anthropic API 金鑰
        """
        self.client = anthropic.Anthropic(api_key=api_key)
        self.presentation = None
        self.using_template = False
        self._layout_cache: Dict[str, Any] = {}
        
    def generate_content_from_outline(self, outline: str, style: str = "professional") -> List[SlideContent]:
        """
        根據大綱生成簡報內容
        
        Args:
            outline: 使用者輸入的簡報大綱
            style: 簡報風格 (professional, creative, minimal, academic)
        
        Returns:
            生成的簡報內容列表
        """
        
        prompt = f"""
        請根據以下大綱生成一份完整的簡報內容。請以 JSON 格式返回，包含每一頁的標題和內容。
        
        簡報大綱：
        {outline}
        
        簡報風格：{style}
        
        請生成包含以下結構的 JSON：
        {{
            "title": "簡報主標題",
            "subtitle": "副標題",
            "slides": [
                {{
                    "title": "頁面標題",
                    "content": ["要點1", "要點2", "要點3"],
                    "layout": "content",
                    "notes": "演講者備註"
                }}
            ]
        }}
        
        要求：
        1. 內容要專業且有深度
        2. 每個要點要簡潔有力
        3. 適當使用數據和例證
        4. 邏輯清晰，層次分明
        5. 根據內容選擇合適的版面配置(content, two_column, image_text)
        """
        
        try:
            response = self.client.messages.create(
                model="claude-3-5-sonnet-20241022",
                max_tokens=4000,
                temperature=0.7,
                messages=[
                    {
                        "role": "user",
                        "content": prompt
                    }
                ]
            )
            
            # 解析 AI 回應
            content = response.content[0].text
            
            # 嘗試提取 JSON 內容
            import re
            json_match = re.search(r'\{.*\}', content, re.DOTALL)
            if json_match:
                json_str = json_match.group()
                presentation_data = json.loads(json_str)
            else:
                raise ValueError("無法從 AI 回應中提取 JSON 內容")
            
            # 轉換為 SlideContent 物件
            slides = []
            
            # 添加標題頁
            slides.append(SlideContent(
                title=presentation_data.get("title", "簡報標題"),
                content=[presentation_data.get("subtitle", "")],
                layout="title"
            ))
            
            # 添加內容頁
            for slide_data in presentation_data.get("slides", []):
                slides.append(SlideContent(
                    title=slide_data.get("title", ""),
                    content=slide_data.get("content", []),
                    layout=slide_data.get("layout", "content"),
                    notes=slide_data.get("notes", "")
                ))
            
            return slides
            
        except Exception as e:
            print(f"生成內容時發生錯誤：{e}")
            return self._generate_fallback_content(outline)
    
    def _generate_fallback_content(self, outline: str) -> List[SlideContent]:
        """當 API 調用失敗時的備用內容生成"""
        lines = outline.strip().split('\n')
        slides = []
        
        # 標題頁
        slides.append(SlideContent(
            title=lines[0] if lines else "簡報標題",
            content=["自動生成的簡報"],
            layout="title"
        ))
        
        # 根據大綱生成內容頁
        current_title = ""
        current_content = []
        
        for line in lines[1:]:
            line = line.strip()
            if not line:
                continue
                
            if line.startswith('#') or line.startswith('*') or line.startswith('-'):
                if current_title and current_content:
                    slides.append(SlideContent(
                        title=current_title,
                        content=current_content,
                        layout="content"
                    ))
                    current_content = []
                current_title = line.lstrip('#*- ')
            else:
                current_content.append(line)
        
        # 添加最後一頁
        if current_title and current_content:
            slides.append(SlideContent(
                title=current_title,
                content=current_content,
                layout="content"
            ))
        
        return slides
    
    def create_presentation(
        self,
        slides: List[SlideContent],
        template: str = "modern",
        template_path: Optional[str] = None,
        add_ending_slide: bool = True
    ) -> Presentation:
        """
        創建 PowerPoint 簡報
        
        Args:
            slides: 簡報內容列表
            template: 簡報模板風格
            template_path: 既有 PPT 模板檔案路徑
            add_ending_slide: 是否添加結尾頁
        
        Returns:
            生成的簡報物件
        """
        self.presentation = self._load_presentation(template_path)
        
        # 設定簡報大小 (16:9)
        if not self.using_template:
            self.presentation.slide_width = Inches(10)
            self.presentation.slide_height = Inches(5.625)
        
        if self.using_template:
            # 使用模板時的特殊處理
            self._create_presentation_from_template(slides, add_ending_slide)
        else:
            # 一般處理
            for slide_content in slides:
                if slide_content.layout == "title":
                    self._add_title_slide(slide_content)
                elif slide_content.layout == "two_column":
                    self._add_two_column_slide(slide_content)
                else:
                    self._add_content_slide(slide_content)
        
        return self.presentation
    
    def _create_presentation_from_template(self, slides: List[SlideContent], add_ending_slide: bool = True):
        """使用模板創建簡報，保留模板的首頁、內容頁、尾頁設計"""
        template_info = getattr(self, 'template_slides_info', {})
        
        # 保存原始投影片的參考
        original_slides = list(self.presentation.slides)
        
        # 清空現有投影片但保留版面配置
        self._clear_template_slides()
        
        # 處理標題頁
        title_slides = [s for s in slides if s.layout == "title"]
        if title_slides:
            for title_slide in title_slides:
                if template_info.get('title_slide') is not None:
                    original_slide_index = template_info['title_slide']
                    if original_slide_index < len(original_slides):
                        layout = original_slides[original_slide_index].slide_layout
                        self._add_slide_with_layout(title_slide, layout, is_title=True)
        
        # 處理內容頁
        content_slides = [s for s in slides if s.layout != "title"]
        content_indices = template_info.get('content_slides', [1])
        
        for i, content_slide in enumerate(content_slides):
            original_slide_index = content_indices[i % len(content_indices)]
            if original_slide_index < len(original_slides):
                layout = original_slides[original_slide_index].slide_layout
                is_two_column = content_slide.layout == "two_column"
                self._add_slide_with_layout(content_slide, layout, is_title=False, is_two_column=is_two_column)
        
        # 添加結尾頁（如果模板中有且用戶要求）
        if add_ending_slide and template_info.get('ending_slide') is not None:
            original_slide_index = template_info['ending_slide']
            if original_slide_index < len(original_slides):
                layout = original_slides[original_slide_index].slide_layout
                ending_content = SlideContent("謝謝聆聽", ["感謝您的聆聽！"], "ending")
                self._add_slide_with_layout(ending_content, layout, is_title=False)
    
    def _clear_template_slides(self):
        """清除模板中的範例內容，但保留版面設計"""
        try:
            slides = self.presentation.slides
            while len(slides) > 0:
                slide_id = slides._sldIdLst[0].rId
                self.presentation.part.drop_rel(slide_id)
                del slides._sldIdLst[0]
        except Exception as exc:
            print(f"清除模板範例投影片時發生錯誤：{exc}")
    

    
    def _add_slide_with_layout(self, content: SlideContent, layout, is_title: bool = False, is_two_column: bool = False):
        """使用指定版面配置添加投影片"""
        slide = self.presentation.slides.add_slide(layout)
        
        # 找到標題佔位符
        title_placeholder = None
        for placeholder in slide.placeholders:
            try:
                if hasattr(placeholder.placeholder_format, 'type'):
                    ph_type = placeholder.placeholder_format.type
                    if ph_type == 1:  # TITLE 類型
                        title_placeholder = placeholder
                        break
            except:
                continue
        
        # 如果沒找到 TITLE 類型，使用第一個文字佔位符
        if not title_placeholder and slide.placeholders:
            for placeholder in slide.placeholders:
                if hasattr(placeholder, 'text_frame'):
                    title_placeholder = placeholder
                    break
        
        # 設定標題
        if title_placeholder and content.title:
            try:
                if hasattr(title_placeholder, 'text_frame'):
                    title_placeholder.text_frame.clear()
                    title_paragraph = title_placeholder.text_frame.paragraphs[0]
                    title_paragraph.text = content.title
                else:
                    title_placeholder.text = content.title
            except Exception as e:
                print(f"設定標題失敗: {e}")
        
        # 找到內容佔位符（排除標題佔位符和頁碼）
        content_placeholders = []
        for placeholder in slide.placeholders:
            if placeholder == title_placeholder:
                continue
            try:
                if hasattr(placeholder.placeholder_format, 'type'):
                    ph_type = placeholder.placeholder_format.type
                    # 排除頁碼 (13) 佔位符
                    if ph_type == 13:  # SLIDE_NUMBER
                        continue
                    if ph_type == 2 and hasattr(placeholder, 'text_frame'):  # BODY 類型
                        content_placeholders.append(placeholder)
            except:
                if hasattr(placeholder, 'text_frame'):
                    content_placeholders.append(placeholder)
        
        # 設定內容
        if is_two_column and len(content_placeholders) >= 1:
            # 雙欄處理
            mid_point = len(content.content) // 2
            left_content = content.content[:mid_point]
            right_content = content.content[mid_point:]
            
            # 左欄
            if content_placeholders:
                tf = content_placeholders[0].text_frame
                tf.clear()
                for i, point in enumerate(left_content):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = point
            
            # 右欄（如果有第二個佔位符）
            if len(content_placeholders) > 1:
                tf = content_placeholders[1].text_frame
                tf.clear()
                for i, point in enumerate(right_content):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = point
        
        elif content_placeholders and content.content:
            # 單欄處理
            tf = content_placeholders[0].text_frame
            tf.clear()
            for i, point in enumerate(content.content):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = point

    def _load_presentation(self, template_path: Optional[str]) -> Presentation:
        """載入既有模板或建立新的簡報"""
        if template_path:
            if os.path.isfile(template_path):
                try:
                    presentation = Presentation(template_path)
                    self.using_template = True
                    self._layout_cache = {}
                    # 保留模板投影片作為參考，不刪除
                    self._analyze_template_slides(presentation)
                    return presentation
                except Exception as exc:
                    print(f"載入模板時發生錯誤：{exc}\n將改用預設樣板。")
            else:
                print("找不到指定的模板檔案，將改用預設樣板。")
        self.using_template = False
        self._layout_cache = {}
        return Presentation()
    
    def _analyze_template_slides(self, presentation: Presentation):
        """分析模板中的投影片結構，識別首頁、內容頁、尾頁"""
        try:
            slides = presentation.slides
            slide_count = len(slides)
            
            self.template_slides_info = {
                'title_slide': None,      # 首頁
                'content_slides': [],     # 內容頁
                'ending_slide': None      # 尾頁
            }
            
            for i, slide in enumerate(slides):
                layout_name = (slide.slide_layout.name or "").lower()
                
                # 簡化的判斷邏輯
                if i == 0:  # 第一張通常是首頁
                    self.template_slides_info['title_slide'] = i
                elif i == slide_count - 1 and slide_count > 2:  # 最後一張可能是尾頁
                    if '標題' in layout_name or 'title' in layout_name:
                        self.template_slides_info['ending_slide'] = i
                    else:
                        self.template_slides_info['content_slides'].append(i)
                else:  # 中間的都是內容頁
                    self.template_slides_info['content_slides'].append(i)
            
            print(f"模板分析結果：首頁={self.template_slides_info['title_slide']}, "
                  f"內容頁={self.template_slides_info['content_slides']}, "
                  f"尾頁={self.template_slides_info['ending_slide']}")
                  
        except Exception as exc:
            print(f"分析模板投影片時發生錯誤：{exc}")
            self.template_slides_info = {'title_slide': None, 'content_slides': [], 'ending_slide': None}
    
    def _identify_slide_type(self, slide) -> str:
        """識別投影片類型"""
        try:
            layout = slide.slide_layout
            layout_name = (layout.name or "").lower()
            
            # 根據版面配置名稱判斷
            if any(keyword in layout_name for keyword in ['title', '標題', '封面', '首頁']):
                # 進一步判斷是否為首頁或尾頁
                slide_index = None
                slides_list = list(slide.part.presentation.slides._sldIdLst)
                for i, slide_id in enumerate(slides_list):
                    if slide_id.rId == slide.part.partname.split('/')[-1].replace('.xml', ''):
                        slide_index = i
                        break
                
                # 如果是第一張，通常是首頁
                if slide_index == 0:
                    return 'title'
                # 如果是最後一張且總數>2，可能是尾頁
                elif slide_index is not None and slide_index == len(slides_list) - 1 and len(slides_list) > 2:
                    return 'ending'
                else:
                    return 'title'  # 其他標題投影片當作標題頁處理
            
            elif 'two' in layout_name or 'object' in layout_name:
                return 'content'
            elif any(keyword in layout_name for keyword in ['ending', '結尾', '謝謝', 'thank', 'end']):
                return 'ending'
            else:
                return 'content'
                
        except Exception:
            return 'content'
    
    def _get_slide_layout(self, layout_type: str):
        """依照需求取得合適的版面配置"""
        if not self.presentation:
            raise ValueError("簡報尚未初始化")
        layout_type = layout_type.lower()

        cache_key = layout_type
        if cache_key in self._layout_cache:
            return self._layout_cache[cache_key]

        layouts = list(self.presentation.slide_layouts)
        if not layouts:
            raise ValueError("模板中沒有可用的版面配置")

        # 如果模板只有少數版面配置，優先使用現有投影片的版面
        if len(layouts) <= 2 and hasattr(self, 'template_slides_info'):
            best_layout = self._get_layout_from_template_slides(layout_type)
            if best_layout:
                self._layout_cache[cache_key] = best_layout
                return best_layout

        best_layout = None
        best_score = float("-inf")

        for layout in layouts:
            score = self._score_layout(layout, layout_type)
            if score > best_score:
                best_layout = layout
                best_score = score

        if not best_layout:
            best_layout = layouts[0]

        self._layout_cache[cache_key] = best_layout
        return best_layout
    
    def _get_layout_from_template_slides(self, layout_type: str):
        """從模板投影片中獲取適合的版面配置"""
        try:
            template_info = getattr(self, 'template_slides_info', {})
            slides = list(self.presentation.slides)
            
            if layout_type == "title" and template_info.get('title_slide') is not None:
                slide_idx = template_info['title_slide']
                if slide_idx < len(slides):
                    return slides[slide_idx].slide_layout
            
            elif layout_type == "two_column":
                # 尋找雙欄版面，優先使用 TWO_OBJECTS
                for slide_idx in template_info.get('content_slides', []):
                    if slide_idx < len(slides):
                        slide = slides[slide_idx]
                        if 'TWO' in slide.slide_layout.name.upper():
                            return slide.slide_layout
                
                # 如果沒找到，使用第一個內容頁版面
                if template_info.get('content_slides'):
                    slide_idx = template_info['content_slides'][0]
                    if slide_idx < len(slides):
                        return slides[slide_idx].slide_layout
            
            else:  # content
                # 使用內容頁版面
                if template_info.get('content_slides'):
                    slide_idx = template_info['content_slides'][0]
                    if slide_idx < len(slides):
                        return slides[slide_idx].slide_layout
            
            return None
            
        except Exception as exc:
            print(f"從模板投影片獲取版面配置時發生錯誤：{exc}")
            return None

    def _score_layout(self, layout, layout_type: str) -> float:
        """給予版面配置評分，以挑選最合適的版型"""
        def placeholder_type(ph) -> Optional[int]:
            fmt = getattr(ph, "placeholder_format", None)
            return fmt.type if fmt is not None else None

        types = [placeholder_type(ph) for ph in layout.placeholders]
        title_count = sum(1 for t in types if t in {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE})
        subtitle_count = sum(1 for t in types if t == PP_PLACEHOLDER.SUBTITLE)
        body_count = sum(1 for t in types if t in {PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT})

        base_score = 0.0
        name = (layout.name or "").lower()

        if layout_type == "title":
            if title_count == 0:
                return float("-inf")
            base_score = 10 * title_count + 3 * subtitle_count - body_count
            if "title" in name or "封面" in name:
                base_score += 5
        elif layout_type == "two_column":
            if body_count < 2:
                return float("-inf")
            base_score = 10 * body_count + 4 * title_count
            if "two" in name or "雙" in name or "比較" in name:
                base_score += 3
        else:  # content
            if body_count == 0:
                return float("-inf")
            base_score = 10 * body_count + 2 * title_count
            if "content" in name or "內文" in name or "章節" in name:
                base_score += 3

        return base_score
    
    def _add_title_slide(self, content: SlideContent):
        """添加標題頁"""
        slide_layout = self._get_slide_layout("title")
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # 設定標題
        title = slide.shapes.title
        if title:
            title.text = content.title
        
        # 設定副標題
        if content.content:
            subtitle = None
            if len(slide.placeholders) > 1:
                subtitle = slide.placeholders[1]
            else:
                for placeholder in slide.placeholders:
                    if placeholder != title and placeholder.has_text_frame:
                        subtitle = placeholder
                        break
            if subtitle and subtitle.has_text_frame:
                subtitle.text = content.content[0] if content.content else ""
        
        # 自訂樣式
        self._apply_modern_style(slide)
    
    def _add_content_slide(self, content: SlideContent):
        """添加內容頁"""
        slide_layout = self._get_slide_layout("content")
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # 設定標題
        title = slide.shapes.title
        if title:
            title.text = content.title
        
        # 設定內容
        body = None
        for placeholder in slide.placeholders:
            if placeholder == title:
                continue
            if getattr(placeholder, "has_text_frame", False):
                body = placeholder
                break
        if not body or not getattr(body, "has_text_frame", False):
            return
        tf = body.text_frame
        tf.clear()
        
        for i, point in enumerate(content.content):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = point
            p.level = 0
            p.font.size = Pt(18)
        
        # 添加備註
        if content.notes:
            slide.notes_slide.notes_text_frame.text = content.notes
        
        # 自訂樣式
        self._apply_modern_style(slide)
    
    def _add_two_column_slide(self, content: SlideContent):
        """添加雙欄版面"""
        slide_layout = self._get_slide_layout("two_column")
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # 設定標題
        title = slide.shapes.title
        if title:
            title.text = content.title
        
        # 分割內容到兩欄
        mid_point = len(content.content) // 2
        left_content = content.content[:mid_point]
        right_content = content.content[mid_point:]
        
        # 左欄
        text_placeholders = [
            placeholder for placeholder in slide.placeholders
            if placeholder != title and getattr(placeholder, "has_text_frame", False)
        ]

        if text_placeholders:
            tf = text_placeholders[0].text_frame
            tf.clear()
            for i, point in enumerate(left_content):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = point
                p.font.size = Pt(16)
        
        # 右欄
        if len(text_placeholders) > 1:
            tf = text_placeholders[1].text_frame
            tf.clear()
            for i, point in enumerate(right_content):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = point
                p.font.size = Pt(16)
        
        self._apply_modern_style(slide)
    
    def _apply_modern_style(self, slide):
        """套用現代化樣式"""
        if self.using_template:
            return
        # 為標題添加樣式
        if slide.shapes.title:
            title = slide.shapes.title
            if title.text_frame and title.text_frame.paragraphs:
                for paragraph in title.text_frame.paragraphs:
                    paragraph.font.bold = True
                    paragraph.font.size = Pt(32)
                    paragraph.font.color.rgb = RGBColor(44, 62, 80)
        
        # 添加裝飾性元素（可選）
        left = Inches(0.5)
        top = Inches(5)
        width = Inches(9)
        height = Inches(0.05)
        
        try:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left, top, width, height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(52, 152, 219)
            shape.line.fill.background()
        except:
            pass  # 忽略樣式錯誤
    
    def save_presentation(self, filename: str):
        """
        儲存簡報檔案
        
        Args:
            filename: 輸出檔案名稱
        """
        if self.presentation:
            self.presentation.save(filename)
            print(f"簡報已儲存至：{filename}")
        else:
            print("尚未生成簡報")
    
    def generate_ppt_code(self, outline: str) -> str:
        """
        生成可以產生 PPT 的 Python 程式碼
        
        Args:
            outline: 簡報大綱
        
        Returns:
            可執行的 Python 程式碼
        """
        slides = self.generate_content_from_outline(outline)
        
        code = f'''
"""
自動生成的 PPT 程式碼
生成時間：{datetime.now().strftime("%Y-%m-%d %H:%M:%S")}
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# 創建簡報
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# 簡報內容
slides_data = {json.dumps([{
    "title": s.title,
    "content": s.content,
    "layout": s.layout,
    "notes": s.notes
} for s in slides], ensure_ascii=False, indent=2)}

# 生成每一頁
for slide_data in slides_data:
    if slide_data["layout"] == "title":
        # 標題頁
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = slide_data["title"]
        if slide_data["content"] and len(slide.placeholders) > 1:
            slide.placeholders[1].text = slide_data["content"][0]
    else:
        # 內容頁
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = slide_data["title"]
        
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()
        
        for i, point in enumerate(slide_data["content"]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = point
            p.font.size = Pt(18)
        
        if slide_data.get("notes"):
            slide.notes_slide.notes_text_frame.text = slide_data["notes"]

# 儲存簡報
prs.save("generated_presentation.pptx")
print("簡報已生成：generated_presentation.pptx")
'''
        return code


def main():
    """主程式"""
    print("="*50)
    print("AI PPT Generator - 類似 Gamma 的自動簡報生成器")
    print("="*50)
    
    # 設定 API 金鑰
    api_key = os.getenv("ANTHROPIC_API_KEY")
    if not api_key:
        api_key = input("請輸入您的 Anthropic API 金鑰: ").strip()
    
    # 創建生成器
    generator = AIpptGenerator(api_key)
    
    # 獲取使用者輸入
    print("\n請輸入您的簡報大綱（輸入 'END' 結束）：")
    outline_lines = []
    while True:
        line = input()
        if line.upper() == 'END':
            break
        outline_lines.append(line)
    
    outline = '\n'.join(outline_lines)
    
    if not outline.strip():
        # 使用範例大綱
        outline = """
        AI 技術的商業應用
        # AI 簡介與發展歷程
        # 主要應用領域
        - 金融科技
        - 醫療健康
        - 零售電商
        - 製造業
        # 成功案例分析
        # 實施策略與挑戰
        # 未來展望
        """
    
    print("\n選擇簡報風格：")
    print("1. Professional (專業)")
    print("2. Creative (創意)")
    print("3. Minimal (極簡)")
    print("4. Academic (學術)")
    
    style_choice = input("請選擇 (1-4，預設為 1): ").strip() or "1"
    styles = {
        "1": "professional",
        "2": "creative", 
        "3": "minimal",
        "4": "academic"
    }
    style = styles.get(style_choice, "professional")
    
    template_path = input("\n若要套用既有 PPT 模板，請輸入檔案路徑（留空表示使用內建樣板）: ").strip()
    if template_path:
        template_path = os.path.expanduser(template_path)
        if not os.path.isfile(template_path):
            print("找不到指定的模板，將改用內建樣板。")
            template_path = ""
    template_path = template_path or None
    
    print("\n正在生成簡報內容...")
    
    try:
        # 生成內容
        slides = generator.generate_content_from_outline(outline, style)
        
        print(f"\n已生成 {len(slides)} 頁簡報")
        
        # 預覽內容
        print("\n簡報預覽：")
        print("-"*40)
        for i, slide in enumerate(slides, 1):
            print(f"第 {i} 頁: {slide.title}")
            for point in slide.content[:3]:  # 顯示前3個要點
                print(f"  • {point}")
            if len(slide.content) > 3:
                print(f"  ... 還有 {len(slide.content)-3} 個要點")
            print()
        
        # 選擇輸出方式
        print("\n請選擇輸出方式：")
        print("1. 直接生成 PPT 檔案")
        print("2. 生成 Python 程式碼")
        print("3. 兩者都要")
        
        output_choice = input("請選擇 (1-3，預設為 1): ").strip() or "1"
        
        if output_choice in ["1", "3"]:
            # 詢問是否添加結尾頁
            add_ending = input("\n是否添加結尾頁？(y/n，預設為 y): ").strip().lower()
            add_ending_slide = add_ending != 'n'
            
            # 創建並儲存簡報
            presentation = generator.create_presentation(slides, template_path=template_path, add_ending_slide=add_ending_slide)
            filename = f"presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
            generator.save_presentation(filename)
        
        if output_choice in ["2", "3"]:
            # 生成程式碼
            code = generator.generate_ppt_code(outline)
            code_filename = f"ppt_generator_{datetime.now().strftime('%Y%m%d_%H%M%S')}.py"
            with open(code_filename, 'w', encoding='utf-8') as f:
                f.write(code)
            print(f"\nPython 程式碼已儲存至：{code_filename}")
            if template_path:
                print("提醒：產出的程式碼採用內建樣板，若需套用模板請於程式中自行載入。")
        
        print("\n生成完成！")
        
    except Exception as e:
        print(f"\n發生錯誤：{e}")
        print("請確認 API 金鑰是否正確，以及是否已安裝必要的套件：")
        print("pip install anthropic python-pptx")


if __name__ == "__main__":
    main()