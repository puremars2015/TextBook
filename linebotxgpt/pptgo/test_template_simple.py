#!/usr/bin/env python3
"""
測試使用模板生成 PPT 的功能 - 簡化版本
"""

import os
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from typing import List, Dict, Optional, Any

class SlideContent:
    """簡報頁面內容結構"""
    def __init__(self, title: str, content: List[str], layout: str = "content", notes: str = ""):
        self.title = title
        self.content = content
        self.layout = layout
        self.notes = notes

class SimpleTemplateGenerator:
    """簡化版的模板 PPT 生成器"""
    
    def __init__(self):
        self.presentation = None
        self.using_template = False
        self.template_slides_info = {}
        self._layout_cache = {}
    
    def _load_presentation(self, template_path: Optional[str]) -> Presentation:
        """載入既有模板或建立新的簡報"""
        if template_path and os.path.isfile(template_path):
            try:
                presentation = Presentation(template_path)
                self.using_template = True
                self._analyze_template_slides(presentation)
                return presentation
            except Exception as exc:
                print(f"載入模板時發生錯誤：{exc}")
        
        self.using_template = False
        return Presentation()
    
    def _analyze_template_slides(self, presentation: Presentation):
        """分析模板中的投影片結構"""
        try:
            slides = presentation.slides
            slide_count = len(slides)
            
            self.template_slides_info = {
                'title_slide': None,      
                'content_slides': [],     
                'ending_slide': None      
            }
            
            for i, slide in enumerate(slides):
                layout_name = (slide.slide_layout.name or "").lower()
                
                if i == 0:  # 第一張通常是首頁
                    self.template_slides_info['title_slide'] = i
                elif i == slide_count - 1 and slide_count > 2:  # 最後一張可能是尾頁
                    if '標題' in layout_name or 'title' in layout_name:
                        self.template_slides_info['ending_slide'] = i
                    else:
                        self.template_slides_info['content_slides'].append(i)
                else:  # 中間的都是內容頁
                    self.template_slides_info['content_slides'].append(i)
            
            print(f"模板分析結果：首頁={self.template_slides_info['title_slide']}, "
                  f"內容頁={self.template_slides_info['content_slides']}, "
                  f"尾頁={self.template_slides_info['ending_slide']}")
                  
        except Exception as exc:
            print(f"分析模板投影片時發生錯誤：{exc}")
    
    def _clear_template_slides(self):
        """清除模板中的範例內容"""
        try:
            slides = self.presentation.slides
            while len(slides) > 0:
                slide_id = slides._sldIdLst[0].rId
                self.presentation.part.drop_rel(slide_id)
                del slides._sldIdLst[0]
        except Exception as exc:
            print(f"清除模板範例投影片時發生錯誤：{exc}")
    
    def _get_layout_by_index(self, slide_index: int):
        """根據投影片索引獲取版面配置"""
        layouts = list(self.presentation.slide_layouts)
        if not layouts:
            return None
        
        # 直接使用第一個（通常也是唯一）版面配置
        return layouts[0]
    
    def create_presentation(self, slides: List[SlideContent], template_path: Optional[str] = None, add_ending_slide: bool = True):
        """創建簡報"""
        self.presentation = self._load_presentation(template_path)
        
        if self.using_template:
            # 不清除現有投影片，而是複製它們的版面配置
            template_info = self.template_slides_info
            
            # 保存原始投影片的參考
            original_slides = list(self.presentation.slides)
            
            # 清除現有投影片
            self._clear_template_slides()
            
            # 添加標題頁
            title_slides = [s for s in slides if s.layout == "title"]
            for title_slide in title_slides:
                if template_info.get('title_slide') is not None:
                    original_slide_index = template_info['title_slide']
                    if original_slide_index < len(original_slides):
                        layout = original_slides[original_slide_index].slide_layout
                        self._add_slide_with_layout(title_slide, layout, is_title=True)
            
            # 添加內容頁
            content_slides = [s for s in slides if s.layout != "title"]
            content_indices = template_info.get('content_slides', [1])
            
            for i, content_slide in enumerate(content_slides):
                original_slide_index = content_indices[i % len(content_indices)]
                if original_slide_index < len(original_slides):
                    layout = original_slides[original_slide_index].slide_layout
                    is_two_column = content_slide.layout == "two_column"
                    self._add_slide_with_layout(content_slide, layout, is_title=False, is_two_column=is_two_column)
            
            # 添加結尾頁
            if add_ending_slide and template_info.get('ending_slide') is not None:
                original_slide_index = template_info['ending_slide']
                if original_slide_index < len(original_slides):
                    layout = original_slides[original_slide_index].slide_layout
                    ending_content = SlideContent("謝謝聆聽", ["Thank you for your attention!"], "ending")
                    self._add_slide_with_layout(ending_content, layout, is_title=False)
        else:
            # 使用預設方式
            for slide_content in slides:
                layout = self.presentation.slide_layouts[0] if slide_content.layout == "title" else self.presentation.slide_layouts[1]
                slide = self.presentation.slides.add_slide(layout)
                
                if slide.shapes.title:
                    slide.shapes.title.text = slide_content.title
                
                # 添加內容
                if len(slide.placeholders) > 1:
                    body = slide.placeholders[1]
                    if hasattr(body, 'text_frame'):
                        tf = body.text_frame
                        tf.clear()
                        for j, point in enumerate(slide_content.content):
                            if j == 0:
                                p = tf.paragraphs[0]
                            else:
                                p = tf.add_paragraph()
                            p.text = point
        
        return self.presentation
    
    def _add_slide_with_layout(self, content: SlideContent, layout, is_title: bool = False, is_two_column: bool = False):
        """使用指定版面配置添加投影片"""
        slide = self.presentation.slides.add_slide(layout)
        print(f"  添加投影片: {content.title}, 版面: {layout.name}, 佔位符數量: {len(slide.placeholders)}")
        
        # 找到標題佔位符
        title_placeholder = None
        for placeholder in slide.placeholders:
            try:
                if hasattr(placeholder.placeholder_format, 'type'):
                    ph_type = placeholder.placeholder_format.type
                    if ph_type == 1:  # TITLE 類型
                        title_placeholder = placeholder
                        break
            except:
                continue
        
        # 如果沒找到 TITLE 類型，使用第一個文字佔位符
        if not title_placeholder and slide.placeholders:
            for placeholder in slide.placeholders:
                if hasattr(placeholder, 'text_frame'):
                    title_placeholder = placeholder
                    break
        
        # 設定標題
        if title_placeholder and content.title:
            try:
                if hasattr(title_placeholder, 'text_frame'):
                    title_placeholder.text_frame.clear()
                    title_paragraph = title_placeholder.text_frame.paragraphs[0]
                    title_paragraph.text = content.title
                    print(f"    ✓ 設定標題: '{content.title}'")
                else:
                    title_placeholder.text = content.title
                    print(f"    ✓ 設定標題 (直接): '{content.title}'")
            except Exception as e:
                print(f"    ✗ 設定標題失敗: {e}")
        
        # 找到內容佔位符（排除標題佔位符和頁碼）
        content_placeholders = []
        for placeholder in slide.placeholders:
            if placeholder == title_placeholder:
                continue
            try:
                if hasattr(placeholder.placeholder_format, 'type'):
                    ph_type = placeholder.placeholder_format.type
                    # 排除頁碼 (13) 佔位符
                    if ph_type == 13:  # SLIDE_NUMBER
                        continue
                    if ph_type == 2 and hasattr(placeholder, 'text_frame'):  # BODY 類型
                        content_placeholders.append(placeholder)
            except:
                if hasattr(placeholder, 'text_frame'):
                    content_placeholders.append(placeholder)
        
        print(f"    找到 {len(content_placeholders)} 個內容佔位符")
        
        # 設定內容
        if is_two_column and len(content_placeholders) >= 1:
            # 雙欄處理
            mid_point = len(content.content) // 2
            left_content = content.content[:mid_point]
            right_content = content.content[mid_point:]
            
            # 左欄
            if content_placeholders:
                tf = content_placeholders[0].text_frame
                tf.clear()
                for i, point in enumerate(left_content):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = point
                print(f"    ✓ 設定左欄內容: {len(left_content)} 項")
            
            # 右欄（如果有第二個佔位符）
            if len(content_placeholders) > 1:
                tf = content_placeholders[1].text_frame
                tf.clear()
                for i, point in enumerate(right_content):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = point
                print(f"    ✓ 設定右欄內容: {len(right_content)} 項")
        
        elif content_placeholders and content.content:
            # 單欄處理
            tf = content_placeholders[0].text_frame
            tf.clear()
            for i, point in enumerate(content.content):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = point
            print(f"    ✓ 設定內容: {len(content.content)} 項")
    

    
    def save_presentation(self, filename: str):
        """儲存簡報"""
        if self.presentation:
            self.presentation.save(filename)
            print(f"簡報已儲存至：{filename}")

def test_template_generation():
    """測試使用模板生成 PPT"""
    
    template_path = "/Users/maenqi/TextBook/linebotxgpt/pptgo/週報20250919馬恩奇.pptx"
    
    generator = SimpleTemplateGenerator()
    
    test_slides = [
        SlideContent(
            title="測試簡報標題",
            content=["使用模板生成的測試簡報", "副標題內容"],
            layout="title"
        ),
        SlideContent(
            title="第一個內容頁",
            content=[
                "這是第一個要點",
                "這是第二個要點", 
                "這是第三個要點",
                "這是第四個要點"
            ],
            layout="content"
        ),
        SlideContent(
            title="雙欄內容頁", 
            content=[
                "左欄第一點",
                "左欄第二點",
                "右欄第一點", 
                "右欄第二點"
            ],
            layout="two_column"
        ),
        SlideContent(
            title="另一個內容頁",
            content=[
                "更多內容要點",
                "詳細說明事項",
                "重要結論總結"
            ],
            layout="content"
        )
    ]
    
    try:
        print("正在使用模板生成簡報...")
        
        presentation = generator.create_presentation(
            slides=test_slides, 
            template_path=template_path,
            add_ending_slide=True
        )
        
        output_filename = f"test_template_output_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        generator.save_presentation(output_filename)
        
        print(f"✓ 成功生成簡報：{output_filename}")
        print(f"✓ 總共生成 {len(presentation.slides)} 張投影片")
        
        print("\n生成的投影片清單:")
        for i, slide in enumerate(presentation.slides):
            layout_name = slide.slide_layout.name
            title_text = ""
            if slide.shapes.title:
                title_text = slide.shapes.title.text
            print(f"投影片 {i+1}: {layout_name} - '{title_text}'")
        
        return True
        
    except Exception as e:
        print(f"✗ 生成簡報時發生錯誤：{e}")
        import traceback
        traceback.print_exc()
        return False

if __name__ == "__main__":
    print("="*50)
    print("測試模板 PPT 生成功能")
    print("="*50)
    
    success = test_template_generation()
    
    if success:
        print("\n🎉 測試成功！模板功能正常運作。")
    else:
        print("\n❌ 測試失敗，請檢查錯誤訊息。")